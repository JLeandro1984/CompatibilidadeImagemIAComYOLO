from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TITLE_COLOR = RGBColor(15, 23, 42)
SUBTITLE_COLOR = RGBColor(51, 65, 85)
ACCENT_COLOR = RGBColor(2, 132, 199)
BODY_COLOR = RGBColor(30, 41, 59)


slides_data = [
    {
        "title": "Visão Geral",
        "subtitle": "ShelfVision MVP",
        "bullets": [
            "Análise de gôndola com IA híbrida: local + backend opcional",
            "Entrega relatório com marcas, contagem, layout e confiança",
            "Opera sem servidor e melhora precisão com YOLO ativo",
        ],
    },
    {
        "title": "Problema de Negócio",
        "subtitle": "Por que isso importa",
        "bullets": [
            "Auditoria manual é lenta e subjetiva",
            "Dificuldade de padronizar leitura por marca",
            "Necessidade de resposta rápida no PDV",
            "Solução gera análise rastreável por evidências",
        ],
    },
    {
        "title": "Arquitetura",
        "subtitle": "Fluxo ponta a ponta",
        "bullets": [
            "Upload da imagem no frontend",
            "Consulta opcional ao backend YOLO",
            "Execução local de COCO-SSD, OCR e heurísticas",
            "Fusão de sinais e geração de relatório final",
        ],
    },
    {
        "title": "IAs Utilizadas",
        "subtitle": "Responsabilidade de cada camada",
        "bullets": [
            "YOLOv8 custom: detecção direta de marcas",
            "COCO-SSD: objetos e embalagens genéricas",
            "Tesseract OCR: leitura de texto e logo",
            "USE + heurísticas: compatibilidade e reforço visual",
        ],
    },
    {
        "title": "Regra de Decisão",
        "subtitle": "Prioridade de confiança",
        "bullets": [
            "1) YOLO acima do threshold",
            "2) OCR + catálogo de marcas",
            "3) Heurísticas de cor e repetição",
            "Sem evidência suficiente: Marca não identificada",
        ],
    },
    {
        "title": "YOLO na Prática",
        "subtitle": "Como funciona no projeto",
        "bullets": [
            "Inferência em única passada na imagem",
            "Saída por objeto: classe, confiança e caixa",
            "Filtro por confidence e IoU (NMS)",
            "API retorna label, confidence e bbox",
        ],
    },
    {
        "title": "Treinamento",
        "subtitle": "Setup atual",
        "bullets": [
            "Dataset em formato YOLO com split train/val",
            "Configuração central em data.yaml",
            "Treino com yolov8n base via train-model.py",
            "Parâmetros atuais: 50 épocas, 640px, patience 20",
        ],
    },
    {
        "title": "Artefatos do Treino",
        "subtitle": "O que é gerado",
        "bullets": [
            "weights best.pt e last.pt",
            "results.csv por época",
            "Curvas PR, F1 e matriz de confusão",
            "best.pt copiado para backend/model para inferência",
        ],
    },
    {
        "title": "Backend",
        "subtitle": "Inferência e operação",
        "bullets": [
            "FastAPI com endpoints health e detect",
            "Modo mock para integração rápida",
            "Modo real com modelo treinado",
            "Configuração por variáveis de ambiente",
        ],
    },
    {
        "title": "Frontend e Relatório",
        "subtitle": "Entrega para o usuário",
        "bullets": [
            "Upload, preview e execução híbrida opcional",
            "Consolidação por fonte YOLO, OCR, Visual e COCO",
            "Resumo com contagem e distribuição na prateleira",
            "Confiança por marca e observações de evidência",
        ],
    },
    {
        "title": "Demo em 3 Minutos",
        "subtitle": "Roteiro recomendado",
        "bullets": [
            "Subir frontend e backend real",
            "Ativar YOLO e testar conexão",
            "Analisar uma imagem e mostrar o relatório",
            "Repetir com backend desligado para mostrar fallback",
        ],
    },
    {
        "title": "Riscos e Limitações",
        "subtitle": "Visão sênior",
        "bullets": [
            "Dataset ainda pequeno para produção",
            "OCR sensível a oclusão e baixa qualidade",
            "Balanceamento de classes ainda crítico",
            "Threshold por marca tende a melhorar precisão",
        ],
    },
    {
        "title": "Próximos Passos",
        "subtitle": "Roadmap técnico",
        "bullets": [
            "Expandir dataset e cenários de captura",
            "Comparar yolov8m e yolov8l em precisão x latência",
            "Versionar experimentos e métricas",
            "Avaliar exportação ONNX ou TensorFlow.js",
        ],
    },
    {
        "title": "Conclusão",
        "subtitle": "Mensagem final",
        "bullets": [
            "MVP funcional já entrega valor",
            "Arquitetura híbrida aumenta resiliência",
            "Multimodelo melhora assertividade",
            "Maior ganho futuro está em escala de dados",
        ],
    },
]


def style_title(shape):
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(38)
            run.font.bold = True
            run.font.color.rgb = TITLE_COLOR


def style_subtitle(shape):
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(20)
            run.font.color.rgb = SUBTITLE_COLOR


def add_divider(slide):
    line = slide.shapes.add_shape(
        1,
        Inches(0.8),
        Inches(1.45),
        Inches(3.0),
        Inches(0.05),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()


def add_executive_slide(prs, title, subtitle, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(12.0), Inches(0.9))
    title_tf = title_box.text_frame
    title_tf.text = title
    style_title(title_box)

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(12.0), Inches(0.6))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = subtitle
    style_subtitle(subtitle_box)

    add_divider(slide)

    body_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.9), Inches(11.6), Inches(4.8))
    body_tf = body_box.text_frame
    body_tf.clear()

    for idx, bullet in enumerate(bullets):
        p = body_tf.paragraphs[0] if idx == 0 else body_tf.add_paragraph()
        p.text = f"• {bullet}"
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(14)
        for run in p.runs:
            run.font.size = Pt(24)
            run.font.color.rgb = BODY_COLOR


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tag_box = slide.shapes.add_textbox(Inches(0.85), Inches(0.65), Inches(3.2), Inches(0.45))
    tag_tf = tag_box.text_frame
    tag_tf.text = "APRESENTAÇÃO TÉCNICA"
    for run in tag_tf.paragraphs[0].runs:
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = ACCENT_COLOR

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.8), Inches(1.8))
    title_tf = title_box.text_frame
    title_tf.text = "ShelfVision MVP"
    for run in title_tf.paragraphs[0].runs:
        run.font.size = Pt(56)
        run.font.bold = True
        run.font.color.rgb = TITLE_COLOR

    subtitle_p = title_tf.add_paragraph()
    subtitle_p.text = "Detecção de Marcas em Gôndola com IA Híbrida"
    for run in subtitle_p.runs:
        run.font.size = Pt(24)
        run.font.color.rgb = SUBTITLE_COLOR

    add_divider(slide)

    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(11.6), Inches(0.6))
    footer_tf = footer_box.text_frame
    footer_tf.text = "Frontend local + Backend YOLO opcional | Relatório orientado a evidências"
    for run in footer_tf.paragraphs[0].runs:
        run.font.size = Pt(16)
        run.font.color.rgb = SUBTITLE_COLOR


def add_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.8), Inches(1.2))
    tf = title_box.text_frame
    tf.text = "Obrigado"
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(56)
        run.font.bold = True
        run.font.color.rgb = TITLE_COLOR

    p2 = tf.add_paragraph()
    p2.text = "Perguntas e próximos passos"
    for run in p2.runs:
        run.font.size = Pt(24)
        run.font.color.rgb = SUBTITLE_COLOR


if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover(prs)

    for item in slides_data:
        add_executive_slide(prs, item["title"], item["subtitle"], item["bullets"])

    add_closing(prs)

    output = "ShelfVision-Apresentacao-v2.pptx"
    prs.save(output)
    print(f"Arquivo gerado: {output}")
