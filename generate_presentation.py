from pptx import Presentation
from pptx.util import Pt


slides_data = [
    {
        "title": "Slide 1 — Visão Geral",
        "bullets": [
            "Projeto: ShelfVision MVP",
            "Objetivo: Identificar marcas em imagens de gôndola com pipeline híbrido (frontend local + backend YOLO opcional)",
            "Resultado final: Relatório automático com marcas, contagem, layout e confiança",
            "Mensagem-chave: Funciona offline no modo local e ganha precisão quando o backend YOLO está ativo.",
        ],
    },
    {
        "title": "Slide 2 — Problema de Negócio",
        "bullets": [
            "Auditoria de gôndola manual é lenta e subjetiva",
            "Difícil padronizar identificação de marca por inspeção humana",
            "Necessidade de resposta rápida para execução no PDV",
            "Valor entregue: análise semi-automatizada e rastreável por evidências de IA.",
        ],
    },
    {
        "title": "Slide 3 — Arquitetura da Solução",
        "bullets": [
            "1) Usuário envia imagem no frontend",
            "2) Frontend opcionalmente consulta POST /detect no backend YOLO",
            "3) Frontend roda modelos locais (COCO-SSD, OCR, heurísticas, semântica)",
            "4) Motor de fusão valida evidências e gera relatório",
            "Fallback: se backend cair, o pipeline local continua.",
        ],
    },
    {
        "title": "Slide 4 — IAs Utilizadas e Função de Cada Uma",
        "bullets": [
            "YOLOv8 custom (Ultralytics/FastAPI): detecção direta de marcas via bbox + confiança",
            "COCO-SSD (TensorFlow.js): detecta objetos genéricos/embalagens para segmentação inicial",
            "Tesseract.js (OCR): lê texto/logo para confirmar marca",
            "Universal Sentence Encoder (TensorFlow.js): mede compatibilidade entre descrição esperada e marcas detectadas",
            "Heurísticas visuais (HSB + regras): reforço por assinatura de cor e padrão de repetição",
            "Mensagem-chave: sistema multimodal de evidências.",
        ],
    },
    {
        "title": "Slide 5 — Fluxo de Decisão (Prioridade)",
        "bullets": [
            "1) YOLO acima do threshold configurado",
            "2) OCR + catálogo de marcas (brands.json)",
            "3) Heurísticas visuais (cor/geometria/repetição)",
            "4) Se não houver evidência suficiente: Marca não identificada",
            "Benefício: redução de falso positivo por múltiplos sinais.",
        ],
    },
    {
        "title": "Slide 6 — YOLO: Como Funciona (didático)",
        "bullets": [
            "YOLO processa a imagem em uma única passada",
            "Cada detecção retorna classe, confiança e bounding box",
            "Pós-processamento aplica confidence threshold e IoU/NMS",
            "No backend a resposta é padronizada em: label, confidence, bbox [x, y, width, height]",
        ],
    },
    {
        "title": "Slide 7 — Dataset e Treinamento",
        "bullets": [
            "Configuração central em data.yaml",
            "Estrutura YOLO: dataset/images/train|val e dataset/labels/train|val",
            "Treino principal via python train-model.py",
            "Base: yolov8n.pt",
            "Hiperparâmetros atuais: epochs=50, imgsz=640, patience=20, device=cpu",
        ],
    },
    {
        "title": "Slide 8 — Arquivos Gerados Após Treinar",
        "bullets": [
            "Pasta principal: runs/detect/brand_detection/",
            "weights/best.pt: melhor checkpoint",
            "weights/last.pt: último checkpoint",
            "results.csv: métricas por época",
            "Gráficos: results.png, PR_curve.png, F1_curve.png, confusion_matrix.png",
            "Pós-treino: best.pt é copiado para backend/model/best.pt",
        ],
    },
    {
        "title": "Slide 9 — Backend de Inferência",
        "bullets": [
            "API FastAPI com GET /health e POST /detect",
            "Configuração por variáveis de ambiente (model path, confidence, iou, mock mode)",
            "Modo mock para teste de integração",
            "Modo YOLO real com modelo treinado",
        ],
    },
    {
        "title": "Slide 10 — Frontend e Relatório Final",
        "bullets": [
            "Upload + preview + detecção híbrida opcional",
            "Consolidação por fonte: YOLO, OCR, Visual, COCO",
            "Relatório com marcas detectadas, contagem e total",
            "Distribuição na prateleira (esquerda/centro/direita)",
            "Confiança por marca e observações por evidência",
        ],
    },
    {
        "title": "Slide 11 — Demo (Roteiro de 3 Minutos)",
        "bullets": [
            "1) Subir frontend: start-frontend.bat",
            "2) Subir backend real: start-backend-real.bat",
            "3) Ativar 'Usar backend YOLO' e testar conexão",
            "4) Enviar imagem e mostrar caixas, resumo e relatório",
            "5) Repetir com backend desligado para mostrar fallback local",
        ],
    },
    {
        "title": "Slide 12 — Pontos de Atenção (visão sênior)",
        "bullets": [
            "Dataset ainda pequeno para alta robustez em produção",
            "Qualidade do OCR cai com baixa resolução/oclusão",
            "Balanceamento de classes e diversidade de cenários impactam mais que só aumentar épocas",
            "Threshold por marca pode melhorar precisão (em vez de threshold único)",
        ],
    },
    {
        "title": "Slide 13 — Próximos Passos",
        "bullets": [
            "Aumentar dataset por marca e diversidade de ambiente",
            "Testar modelos maiores (yolov8m/yolov8l) para precisão x latência",
            "Versionar experimentos e métricas (MLOps básico)",
            "Exportar para ONNX/TensorFlow.js para reduzir dependência de backend",
        ],
    },
    {
        "title": "Slide 14 — Conclusão",
        "bullets": [
            "Projeto já entrega valor como MVP funcional",
            "Arquitetura híbrida reduz risco operacional (fallback local)",
            "Estratégia multimodelo aumenta qualidade de decisão",
            "Próximo ganho relevante: escala de dataset + governança de treino",
        ],
    },
]


def format_text_frame(text_frame, bullets):
    text_frame.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(20)


def build_presentation(output_path: str):
    prs = Presentation()

    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = "ShelfVision MVP"
    title_slide.placeholders[1].text = "Detecção de Marcas em Gôndola\nApresentação técnica e executiva"

    for item in slides_data:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = item["title"]
        body = slide.shapes.placeholders[1].text_frame
        format_text_frame(body, item["bullets"])

    prs.save(output_path)


if __name__ == "__main__":
    build_presentation("ShelfVision-Apresentacao.pptx")
    print("Arquivo gerado: ShelfVision-Apresentacao.pptx")
